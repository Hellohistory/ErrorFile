import bz2
import gzip
import json
import lzma
import shutil
import sqlite3
import tarfile
import tempfile
import wave
import unittest
import zipfile
from pathlib import Path

import py7zr
from ErrorFile import inspect_file, inspect_file_report, inspect_files
from ErrorFile.report import TAG_INVALID_MODE, TAG_NOT_FOUND, TAG_OK
from PIL import Image
from PyPDF2 import PdfWriter
from docx import Document
from openpyxl import Workbook
from pptx import Presentation
from tests.good_xls_bytes import GOOD_XLS_BYTES
from tests.sample_binary_assets import (
    GOOD_FLAC_BYTES,
    GOOD_MP3_BYTES,
    GOOD_MP4_BYTES,
    GOOD_OGG_BYTES,
    GOOD_RAR_BYTES,
)

SVG_SAMPLE = (
    "<svg xmlns='http://www.w3.org/2000/svg' width='10' height='10'>"
    "<rect width='10' height='10' fill='red'/></svg>"
)


class TestFileInspector(unittest.TestCase):
    @classmethod
    def setUpClass(cls):
        cls.temp_dir = tempfile.mkdtemp(prefix="errorfile-tests-")
        cls.good_files = {}
        cls.bad_files = {}
        cls.image_extensions = [
            ".jpg",
            ".jpeg",
            ".png",
            ".gif",
            ".bmp",
            ".webp",
            ".tiff",
            ".svg",
        ]
        cls._prepare_image_files()
        cls._prepare_pdf_files()
        cls._prepare_excel_files()
        cls._prepare_docx_files()
        cls._prepare_pptx_files()
        cls._prepare_zip_files()
        cls._prepare_rar_files()
        cls._prepare_compressed_files()
        cls._prepare_text_files()
        cls._prepare_media_files()

    @classmethod
    def tearDownClass(cls):
        shutil.rmtree(cls.temp_dir, ignore_errors=True)

    @classmethod
    def _register(cls, ext, good_path, bad_path):
        cls.good_files[ext] = str(good_path)
        cls.bad_files[ext] = str(bad_path)

    @classmethod
    def _prepare_image_files(cls):
        format_mapping = {
            ".jpg": "JPEG",
            ".jpeg": "JPEG",
            ".png": "PNG",
            ".gif": "GIF",
            ".bmp": "BMP",
            ".webp": "WEBP",
            ".tiff": "TIFF",
        }
        for ext in cls.image_extensions:
            good_path = Path(cls.temp_dir) / f"good_image{ext}"
            bad_path = Path(cls.temp_dir) / f"bad_image{ext}"
            if ext == ".svg":
                good_path.write_text(SVG_SAMPLE, encoding="utf-8")
                bad_path.write_text("not an svg", encoding="utf-8")
            else:
                image = Image.new("RGB", (10, 10), color="red")
                fmt = format_mapping[ext]
                if fmt == "GIF":
                    image = image.convert("P")
                image.save(good_path, format=fmt)
                bad_path.write_bytes(b"not a valid image file")
            cls._register(ext, good_path, bad_path)

    @classmethod
    def _prepare_pdf_files(cls):
        good_path = Path(cls.temp_dir) / "good.pdf"
        bad_path = Path(cls.temp_dir) / "bad.pdf"
        writer = PdfWriter()
        writer.add_blank_page(width=72, height=72)
        with open(good_path, "wb") as file:
            writer.write(file)
        bad_path.write_text("not a pdf", encoding="utf-8")
        cls._register(".pdf", good_path, bad_path)

    @classmethod
    def _prepare_excel_files(cls):
        good_xlsx = Path(cls.temp_dir) / "good.xlsx"
        bad_xlsx = Path(cls.temp_dir) / "bad.xlsx"
        workbook = Workbook()
        sheet = workbook.active
        sheet["A1"] = "ok"
        workbook.save(good_xlsx)
        bad_xlsx.write_text("not an xlsx", encoding="utf-8")
        cls._register(".xlsx", good_xlsx, bad_xlsx)

        good_xls = Path(cls.temp_dir) / "good.xls"
        bad_xls = Path(cls.temp_dir) / "bad.xls"
        good_xls.write_bytes(GOOD_XLS_BYTES)
        bad_xls.write_text("not an xls", encoding="utf-8")
        cls._register(".xls", good_xls, bad_xls)

    @classmethod
    def _prepare_docx_files(cls):
        good_path = Path(cls.temp_dir) / "good.docx"
        bad_path = Path(cls.temp_dir) / "bad.docx"
        document = Document()
        document.add_paragraph("hello world")
        document.save(good_path)
        bad_path.write_text("not a docx", encoding="utf-8")
        cls._register(".docx", good_path, bad_path)

    @classmethod
    def _prepare_pptx_files(cls):
        good_path = Path(cls.temp_dir) / "good.pptx"
        bad_path = Path(cls.temp_dir) / "bad.pptx"
        presentation = Presentation()
        slide_layout = presentation.slide_layouts[5]
        presentation.slides.add_slide(slide_layout)
        presentation.save(good_path)
        bad_path.write_text("not a pptx", encoding="utf-8")
        cls._register(".pptx", good_path, bad_path)

    @classmethod
    def _prepare_zip_files(cls):
        good_path = Path(cls.temp_dir) / "good.zip"
        bad_path = Path(cls.temp_dir) / "bad.zip"
        with zipfile.ZipFile(good_path, "w") as archive:
            archive.writestr("hello.txt", "zip ok")
        bad_path.write_text("not a zip", encoding="utf-8")
        cls._register(".zip", good_path, bad_path)

    @classmethod
    def _prepare_rar_files(cls):
        good_path = Path(cls.temp_dir) / "good.rar"
        bad_path = Path(cls.temp_dir) / "bad.rar"
        good_path.write_bytes(GOOD_RAR_BYTES)
        bad_path.write_text("not a rar", encoding="utf-8")
        cls._register(".rar", good_path, bad_path)

    @classmethod
    def _prepare_compressed_files(cls):
        gzip_good = Path(cls.temp_dir) / "good.gz"
        gzip_bad = Path(cls.temp_dir) / "bad.gz"
        with gzip.open(gzip_good, "wb") as gzip_file:
            gzip_file.write(b"gzip ok")
        gzip_bad.write_text("not a gzip", encoding="utf-8")
        cls._register(".gz", gzip_good, gzip_bad)

        bzip_good = Path(cls.temp_dir) / "good.bz2"
        bzip_bad = Path(cls.temp_dir) / "bad.bz2"
        with bz2.BZ2File(bzip_good, "wb") as bzip_file:
            bzip_file.write(b"bz2 ok")
        bzip_bad.write_text("not a bz2", encoding="utf-8")
        cls._register(".bz2", bzip_good, bzip_bad)

        seven_zip_good = Path(cls.temp_dir) / "good.7z"
        seven_zip_bad = Path(cls.temp_dir) / "bad.7z"
        payload_txt = Path(cls.temp_dir) / "archive_payload.txt"
        payload_txt.write_text("archive ok", encoding="utf-8")
        with py7zr.SevenZipFile(seven_zip_good, "w") as archive:
            archive.write(payload_txt, arcname="payload.txt")
        seven_zip_bad.write_text("not a 7z", encoding="utf-8")
        cls._register(".7z", seven_zip_good, seven_zip_bad)

        tar_payload = Path(cls.temp_dir) / "tar_payload.txt"
        tar_payload.write_text("tar ok", encoding="utf-8")

        tar_good = Path(cls.temp_dir) / "good.tar"
        tar_bad = Path(cls.temp_dir) / "bad.tar"
        cls._create_tar_archive(tar_good, "w", tar_payload)
        tar_bad.write_text("not a tar", encoding="utf-8")
        cls._register(".tar", tar_good, tar_bad)

        tar_gz_good = Path(cls.temp_dir) / "good.tar.gz"
        tar_gz_bad = Path(cls.temp_dir) / "bad.tar.gz"
        cls._create_tar_archive(tar_gz_good, "w:gz", tar_payload)
        tar_gz_bad.write_text("not a tar.gz", encoding="utf-8")
        cls._register(".tar.gz", tar_gz_good, tar_gz_bad)

        tar_bz2_good = Path(cls.temp_dir) / "good.tar.bz2"
        tar_bz2_bad = Path(cls.temp_dir) / "bad.tar.bz2"
        cls._create_tar_archive(tar_bz2_good, "w:bz2", tar_payload)
        tar_bz2_bad.write_text("not a tar.bz2", encoding="utf-8")
        cls._register(".tar.bz2", tar_bz2_good, tar_bz2_bad)

        tar_xz_good = Path(cls.temp_dir) / "good.tar.xz"
        tar_xz_bad = Path(cls.temp_dir) / "bad.tar.xz"
        cls._create_tar_archive(tar_xz_good, "w:xz", tar_payload)
        tar_xz_bad.write_text("not a tar.xz", encoding="utf-8")
        cls._register(".tar.xz", tar_xz_good, tar_xz_bad)

        xz_good = Path(cls.temp_dir) / "good.xz"
        xz_bad = Path(cls.temp_dir) / "bad.xz"
        with lzma.open(xz_good, "wb") as xz_file:
            xz_file.write(b"xz ok")
        xz_bad.write_text("not an xz", encoding="utf-8")
        cls._register(".xz", xz_good, xz_bad)

    @classmethod
    def _create_tar_archive(cls, path: Path, mode: str, payload: Path) -> None:
        with tarfile.open(path, mode) as tar:
            tar.add(payload, arcname="payload.txt")

    @classmethod
    def _prepare_media_files(cls):
        mp3_path = Path(cls.temp_dir) / "good.mp3"
        bad_mp3 = Path(cls.temp_dir) / "bad.mp3"
        mp3_path.write_bytes(GOOD_MP3_BYTES)
        bad_mp3.write_text("not an mp3", encoding="utf-8")
        cls._register(".mp3", mp3_path, bad_mp3)

        mp4_path = Path(cls.temp_dir) / "good.mp4"
        bad_mp4 = Path(cls.temp_dir) / "bad.mp4"
        mp4_path.write_bytes(GOOD_MP4_BYTES)
        bad_mp4.write_text("not an mp4", encoding="utf-8")
        cls._register(".mp4", mp4_path, bad_mp4)

        flac_path = Path(cls.temp_dir) / "good.flac"
        bad_flac = Path(cls.temp_dir) / "bad.flac"
        flac_path.write_bytes(GOOD_FLAC_BYTES)
        bad_flac.write_text("not a flac", encoding="utf-8")
        cls._register(".flac", flac_path, bad_flac)

        ogg_path = Path(cls.temp_dir) / "good.ogg"
        bad_ogg = Path(cls.temp_dir) / "bad.ogg"
        ogg_path.write_bytes(GOOD_OGG_BYTES)
        bad_ogg.write_text("not an ogg", encoding="utf-8")
        cls._register(".ogg", ogg_path, bad_ogg)

        wav_path = Path(cls.temp_dir) / "good.wav"
        bad_wav = Path(cls.temp_dir) / "bad.wav"
        with wave.open(str(wav_path), "wb") as wav_file:
            wav_file.setnchannels(1)
            wav_file.setsampwidth(2)
            wav_file.setframerate(8000)
            wav_file.writeframes(b"\x00\x00" * 800)
        bad_wav.write_text("not a wav", encoding="utf-8")
        cls._register(".wav", wav_path, bad_wav)

    @classmethod
    def _prepare_text_files(cls):
        json_good = Path(cls.temp_dir) / "good.json"
        json_bad = Path(cls.temp_dir) / "bad.json"
        json_good.write_text(json.dumps({"status": "ok", "value": 1}), encoding="utf-8")
        json_bad.write_text('{"status": }', encoding="utf-8")
        cls._register(".json", json_good, json_bad)

        xml_good = Path(cls.temp_dir) / "good.xml"
        xml_bad = Path(cls.temp_dir) / "bad.xml"
        xml_good.write_text("<root><item>ok</item></root>", encoding="utf-8")
        xml_bad.write_text("<root>", encoding="utf-8")
        cls._register(".xml", xml_good, xml_bad)

        text_good = Path(cls.temp_dir) / "good.txt"
        text_bad = Path(cls.temp_dir) / "bad.txt"
        text_good.write_text("hello text", encoding="utf-8")
        text_bad.write_bytes(b"\x00\x00\xff\xfe")
        cls._register(".txt", text_good, text_bad)

        csv_good = Path(cls.temp_dir) / "good.csv"
        csv_bad = Path(cls.temp_dir) / "bad.csv"
        csv_good.write_text("name,age\nalice,30\n", encoding="utf-8")
        csv_bad.write_text('name,age\n"alice,30\n', encoding="utf-8")
        cls._register(".csv", csv_good, csv_bad)

        html_good = Path(cls.temp_dir) / "good.html"
        html_bad = Path(cls.temp_dir) / "bad.html"
        html_good.write_text("<html><body><p>ok</p></body></html>", encoding="utf-8")
        html_bad.write_text("plain text without tags", encoding="utf-8")
        cls._register(".html", html_good, html_bad)

        ini_good = Path(cls.temp_dir) / "good.ini"
        ini_bad = Path(cls.temp_dir) / "bad.ini"
        ini_good.write_text("[main]\nname=ok\n", encoding="utf-8")
        ini_bad.write_text("name=missing-section\n", encoding="utf-8")
        cls._register(".ini", ini_good, ini_bad)

        ndjson_good = Path(cls.temp_dir) / "good.ndjson"
        ndjson_bad = Path(cls.temp_dir) / "bad.ndjson"
        ndjson_good.write_text('{"name":"alice"}\n{"name":"bob"}\n', encoding="utf-8")
        ndjson_bad.write_text('{"name":"alice"}\n{"name":}\n', encoding="utf-8")
        cls._register(".ndjson", ndjson_good, ndjson_bad)

        tsv_good = Path(cls.temp_dir) / "good.tsv"
        tsv_bad = Path(cls.temp_dir) / "bad.tsv"
        tsv_good.write_text("name\tage\nalice\t30\n", encoding="utf-8")
        tsv_bad.write_text('name\tage\n"alice\t30\n', encoding="utf-8")
        cls._register(".tsv", tsv_good, tsv_bad)

        rtf_good = Path(cls.temp_dir) / "good.rtf"
        rtf_bad = Path(cls.temp_dir) / "bad.rtf"
        rtf_good.write_text("{\\rtf1\\ansi This is rtf}", encoding="utf-8")
        rtf_bad.write_text("This is not rtf", encoding="utf-8")
        cls._register(".rtf", rtf_good, rtf_bad)

        eml_good = Path(cls.temp_dir) / "good.eml"
        eml_bad = Path(cls.temp_dir) / "bad.eml"
        eml_good.write_text(
            "From: test@example.com\n"
            "To: you@example.com\n"
            "Subject: test\n"
            "\n"
            "hello\n",
            encoding="utf-8",
        )
        eml_bad.write_text("body only no headers", encoding="utf-8")
        cls._register(".eml", eml_good, eml_bad)

        toml_good = Path(cls.temp_dir) / "good.toml"
        toml_bad = Path(cls.temp_dir) / "bad.toml"
        toml_good.write_text(
            'title = "Example"\n'
            "[owner]\n"
            'name = "alice"\n',
            encoding="utf-8",
        )
        toml_bad.write_text('title = "Example"\nowner = [\n', encoding="utf-8")
        cls._register(".toml", toml_good, toml_bad)

        yaml_good = Path(cls.temp_dir) / "good.yaml"
        yaml_bad = Path(cls.temp_dir) / "bad.yaml"
        yaml_good.write_text("name: alice\nage: 30\n", encoding="utf-8")
        yaml_bad.write_text("name alice\nage 30\n", encoding="utf-8")
        cls._register(".yaml", yaml_good, yaml_bad)

        yml_good = Path(cls.temp_dir) / "good.yml"
        yml_bad = Path(cls.temp_dir) / "bad.yml"
        yml_good.write_text("- one\n- two\n", encoding="utf-8")
        yml_bad.write_text("just plain words", encoding="utf-8")
        cls._register(".yml", yml_good, yml_bad)

        msg_good = Path(cls.temp_dir) / "good.msg"
        msg_bad = Path(cls.temp_dir) / "bad.msg"
        msg_good.write_bytes(b"\xD0\xCF\x11\xE0\xA1\xB1\x1A\xE1" + b"\x00" * 64)
        msg_bad.write_text("not a msg", encoding="utf-8")
        cls._register(".msg", msg_good, msg_bad)

        sqlite_good = Path(cls.temp_dir) / "good.sqlite"
        sqlite_bad = Path(cls.temp_dir) / "bad.sqlite"
        with sqlite3.connect(str(sqlite_good)) as conn:
            conn.execute("CREATE TABLE sample (id INTEGER PRIMARY KEY, name TEXT)")
            conn.execute("INSERT INTO sample(name) VALUES ('ok')")
            conn.commit()
        sqlite_bad.write_text("not a sqlite db", encoding="utf-8")
        cls._register(".sqlite", sqlite_good, sqlite_bad)

        db_good = Path(cls.temp_dir) / "good.db"
        db_bad = Path(cls.temp_dir) / "bad.db"
        with sqlite3.connect(str(db_good)) as conn:
            conn.execute("CREATE TABLE sample (id INTEGER PRIMARY KEY, value TEXT)")
            conn.execute("INSERT INTO sample(value) VALUES ('ok')")
            conn.commit()
        db_bad.write_text("not a db", encoding="utf-8")
        cls._register(".db", db_good, db_bad)

    def test_good_images(self):
        for ext in self.image_extensions:
            with self.subTest(ext=ext):
                is_ok, message = inspect_file(self.good_files[ext], mode="deep")
                self.assertTrue(is_ok, f"{ext} image should pass: {message}")

    def test_corrupted_images(self):
        for ext in self.image_extensions:
            with self.subTest(ext=ext):
                is_ok, message = inspect_file(self.bad_files[ext], mode="deep")
                self.assertFalse(is_ok, f"{ext} corrupted image should fail: {message}")

    def test_pdf_files(self):
        good_path = self.good_files[".pdf"]
        bad_path = self.bad_files[".pdf"]
        is_ok, message = inspect_file(good_path)
        self.assertTrue(is_ok, f"PDF should pass: {message}")
        is_ok, message = inspect_file(bad_path)
        self.assertFalse(is_ok, "Corrupted PDF should fail.")

    def test_excel_files(self):
        for ext in (".xlsx", ".xls"):
            with self.subTest(ext=ext):
                is_ok, message = inspect_file(self.good_files[ext])
                self.assertTrue(is_ok, f"{ext} should pass: {message}")
                is_ok, message = inspect_file(self.bad_files[ext])
                self.assertFalse(is_ok, f"Corrupted {ext} should fail.")

    def test_docx_files(self):
        is_ok, message = inspect_file(self.good_files[".docx"])
        self.assertTrue(is_ok, f"DOCX should pass: {message}")
        is_ok, message = inspect_file(self.bad_files[".docx"])
        self.assertFalse(is_ok, "Corrupted DOCX should fail.")

    def test_pptx_files(self):
        is_ok, message = inspect_file(self.good_files[".pptx"])
        self.assertTrue(is_ok, f"PPTX should pass: {message}")
        is_ok, message = inspect_file(self.bad_files[".pptx"])
        self.assertFalse(is_ok, "Corrupted PPTX should fail.")

    def test_zip_files(self):
        is_ok, message = inspect_file(self.good_files[".zip"])
        self.assertTrue(is_ok, f"ZIP should pass: {message}")
        is_ok, message = inspect_file(self.bad_files[".zip"])
        self.assertFalse(is_ok, "Corrupted ZIP should fail.")

    def test_rar_files(self):
        is_ok, message = inspect_file(self.good_files[".rar"])
        self.assertTrue(is_ok, f"RAR should pass: {message}")
        is_ok, message = inspect_file(self.bad_files[".rar"])
        self.assertFalse(is_ok, "Corrupted RAR should fail.")

    def test_compressed_files(self):
        for ext in (".gz", ".bz2", ".xz", ".tar", ".tar.gz", ".tar.bz2", ".tar.xz", ".7z"):
            with self.subTest(ext=ext):
                is_ok, message = inspect_file(self.good_files[ext])
                self.assertTrue(is_ok, f"{ext} should pass: {message}")
                is_ok, message = inspect_file(self.bad_files[ext])
                self.assertFalse(is_ok, f"Corrupted {ext} should fail.")

    def test_text_files(self):
        for ext in (
            ".json",
            ".ndjson",
            ".xml",
            ".txt",
            ".csv",
            ".tsv",
            ".html",
            ".ini",
            ".rtf",
            ".eml",
            ".toml",
            ".yaml",
            ".yml",
            ".msg",
            ".sqlite",
            ".db",
        ):
            with self.subTest(ext=ext):
                is_ok, message = inspect_file(self.good_files[ext])
                self.assertTrue(is_ok, f"{ext} should pass: {message}")
                is_ok, message = inspect_file(self.bad_files[ext])
                self.assertFalse(is_ok, f"Corrupted {ext} should fail.")

    def test_media_files(self):
        for ext in (".mp3", ".mp4", ".flac", ".ogg", ".wav"):
            with self.subTest(ext=ext):
                is_ok, message = inspect_file(self.good_files[ext])
                self.assertTrue(is_ok, f"{ext} should pass: {message}")
                is_ok, message = inspect_file(self.bad_files[ext])
                self.assertFalse(is_ok, f"Corrupted {ext} should fail.")

    def test_fast_mode(self):
        report = inspect_file_report(self.good_files[".jpg"], mode="fast")
        self.assertTrue(report.ok)
        self.assertIn(TAG_OK, report.tags)

    def test_invalid_mode(self):
        report = inspect_file_report(self.good_files[".jpg"], mode="unknown")
        self.assertFalse(report.ok)
        self.assertIn(TAG_INVALID_MODE, report.tags)

    def test_non_existent_file(self):
        report = inspect_file_report("non_existent_file_12345.xyz")
        self.assertFalse(report.ok)
        self.assertIn(TAG_NOT_FOUND, report.tags)

    def test_signature_precheck(self):
        fake_pdf = Path(self.temp_dir) / "fake_header.pdf"
        fake_pdf.write_text("this is not a real pdf", encoding="utf-8")
        report = inspect_file_report(str(fake_pdf), mode="deep")
        self.assertFalse(report.ok)
        self.assertIn("invalid_format", report.tags)

    def test_inspect_files_deduplicate_paths(self):
        path = self.good_files[".png"]
        reports = inspect_files([path, path, path], mode="deep")
        self.assertEqual(len(reports), 3)
        for report in reports:
            self.assertTrue(report.ok)
            self.assertIn(TAG_OK, report.tags)

    def test_signature_precheck_denylist_disables_header_check(self):
        fake_pdf = Path(self.temp_dir) / "fake_header_denylist.pdf"
        fake_pdf.write_text("this is not a real pdf", encoding="utf-8")

        enabled = inspect_file_report(str(fake_pdf), mode="deep")
        disabled = inspect_file_report(
            str(fake_pdf),
            mode="deep",
            signature_precheck_denylist=[".pdf"],
        )

        self.assertFalse(enabled.ok)
        self.assertFalse(disabled.ok)
        self.assertIn("signature mismatch", enabled.message.lower())
        self.assertNotIn("signature mismatch", disabled.message.lower())

    def test_staged_deep_preserves_deep_mode(self):
        good = self.good_files[".png"]
        bad = self.bad_files[".pdf"]
        reports = inspect_files([good, bad], mode="deep", staged_deep=True)
        self.assertEqual(2, len(reports))
        for report in reports:
            self.assertEqual("deep", report.mode)


if __name__ == "__main__":
    unittest.main()
