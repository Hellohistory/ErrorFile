# ErrorFile/Detection/PowerPointInspector.py
"""PowerPoint inspection utilities."""

import zipfile

from pptx import Presentation
from pptx.exc import PackageNotFoundError

from ..report import (
    TAG_CORRUPTED,
    TAG_INVALID_FORMAT,
    TAG_IO_ERROR,
    fail_finding,
    ok_finding,
)


def _fast_check_pptx(file_path):
    try:
        with zipfile.ZipFile(file_path, "r") as archive:
            names = set(archive.namelist())
            required = {"[Content_Types].xml", "ppt/presentation.xml"}
            if not required.issubset(names):
                return fail_finding(
                    "PPTX missing required parts; file may be corrupted.",
                    TAG_CORRUPTED,
                    TAG_INVALID_FORMAT,
                )
    except zipfile.BadZipFile as exc:
        return fail_finding(
            f"PPTX is not a valid ZIP container: {exc}",
            TAG_INVALID_FORMAT,
            error=str(exc),
        )
    except Exception as exc:
        return fail_finding(
            f"PPTX fast check failed: {exc}",
            TAG_IO_ERROR,
            error=str(exc),
        )
    return ok_finding("PPTX fast check passed.")


def check_pptx_file(file_path, mode="deep"):
    """Check .pptx integrity."""
    if mode == "fast":
        return _fast_check_pptx(file_path)
    try:
        presentation = Presentation(file_path)
        for slide in presentation.slides:
            for shape in slide.shapes:
                _ = shape.has_text_frame
        return ok_finding("PPTX deep check passed.")
    except PackageNotFoundError as exc:
        return fail_finding(
            f"PPTX corrupted or invalid: {exc}",
            TAG_CORRUPTED,
            TAG_INVALID_FORMAT,
            error=str(exc),
        )
    except Exception as exc:  # pragma: no cover - depends on third-party lib
        return fail_finding(
            f"PPTX deep check failed: {exc}",
            TAG_IO_ERROR,
            error=str(exc),
        )
